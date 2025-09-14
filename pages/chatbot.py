"""
Simplified Chatbot Page - Clean, Minimal Implementation
"""

import streamlit as st
import logging
from datetime import datetime

# Configure logging
logger = logging.getLogger(__name__)

def run_async(coro):
    """Run async function in Streamlit context."""
    import asyncio
    try:
        loop = asyncio.get_event_loop()
    except RuntimeError:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
    
    try:
        return loop.run_until_complete(coro)
    except Exception as e:
        logger.error(f"Error running async operation: {e}")
        raise

def save_conversation():
    """Save current conversation to database."""
    try:
        if not st.session_state.get('authenticated') or not st.session_state.get('user_id'):
            return False
        
        if not st.session_state.chat_messages:
            return True  # Nothing to save
        
        logger.info("Saving conversation to database...")
        
        # Import services
        from services.conversation_service import conversation_service
        from services.user_service import user_service
        
        # Get user UUID
        user_data = run_async(user_service.get_user_by_id(st.session_state.user_id))
        if not user_data:
            logger.error("User not found for conversation saving")
            return False
        
        # Create or update conversation
        if not st.session_state.current_conversation_id:
            # Create new conversation
            title = generate_conversation_title()
            conversation_id = run_async(conversation_service.create_conversation(
                user_data['id'], 
                title, 
                st.session_state.get('selected_model_mode', 'normal')
            ))
            st.session_state.current_conversation_id = conversation_id
            logger.info(f"✅ Created new conversation: {conversation_id}")
        
        # Update conversation with current messages
        success = run_async(conversation_service.update_conversation(
            user_data['id'],
            st.session_state.current_conversation_id,
            {'messages': st.session_state.chat_messages}
        ))
        
        if success:
            logger.info(f"✅ Conversation saved: {len(st.session_state.chat_messages)} messages")
            return True
        else:
            logger.error("Failed to save conversation")
            return False
            
    except Exception as e:
        logger.error(f"Error saving conversation: {e}")
        return False

def generate_conversation_title():
    """Generate a title for the conversation based on the first user message."""
    if st.session_state.chat_messages:
        first_user_msg = next((msg for msg in st.session_state.chat_messages if msg["role"] == "user"), None)
        if first_user_msg:
            content = first_user_msg["content"]
            # Create title from first 50 characters
            title = content[:50] + "..." if len(content) > 50 else content
            return title
    
    # Fallback title
    return f"Chat {datetime.now().strftime('%m/%d %H:%M')}"

def load_conversation_messages():
    """Load messages for the current conversation."""
    try:
        current_conv_id = st.session_state.get('current_conversation_id')
        if current_conv_id and st.session_state.get('conversations'):
            conversation = st.session_state.conversations.get(current_conv_id)
            if conversation:
                st.session_state.chat_messages = conversation.get('messages', [])
                logger.info(f"Loaded {len(st.session_state.chat_messages)} messages for conversation {current_conv_id}")
            else:
                st.session_state.chat_messages = []
        else:
            st.session_state.chat_messages = []
    except Exception as e:
        logger.error(f"Error loading conversation messages: {e}")
        st.session_state.chat_messages = []

def process_uploaded_document(uploaded_file):
    """Process uploaded document and return text content with enhanced support for images and PPTX."""
    try:
        # Read file content
        file_content = uploaded_file.read()
        uploaded_file.seek(0)  # Reset file pointer
        
        # Process based on file type
        if uploaded_file.type == "text/plain" or uploaded_file.name.endswith('.txt'):
            text_content = file_content.decode('utf-8')
            
        elif uploaded_file.name.endswith('.md'):
            text_content = file_content.decode('utf-8')
            
        elif uploaded_file.type == "application/pdf" or uploaded_file.name.endswith('.pdf'):
            try:
                import PyPDF2
                import io
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                text_content = ""
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content += f"\n--- Page {page_num} ---\n{page_text}\n"
                
                if not text_content.strip():
                    text_content = f"[PDF Document: {uploaded_file.name} - {len(file_content)} bytes - No extractable text found]"
                    
            except Exception as pdf_error:
                logger.warning(f"PDF processing failed: {pdf_error}")
                text_content = f"[PDF Document: {uploaded_file.name} - {len(file_content)} bytes - Text extraction failed: {pdf_error}]"
                
        elif uploaded_file.name.endswith('.docx'):
            try:
                import docx2txt
                import tempfile
                import os
                
                with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
                    tmp_file.write(file_content)
                    tmp_path = tmp_file.name
                
                try:
                    text_content = docx2txt.process(tmp_path)
                    if not text_content.strip():
                        # Try alternative method with python-docx
                        try:
                            from docx import Document
                            doc = Document(tmp_path)
                            paragraphs = []
                            for paragraph in doc.paragraphs:
                                if paragraph.text.strip():
                                    paragraphs.append(paragraph.text)
                            text_content = "\n".join(paragraphs)
                        except:
                            raise Exception("No text extracted with any method")
                finally:
                    os.unlink(tmp_path)
                    
            except Exception as docx_error:
                logger.warning(f"DOCX processing failed: {docx_error}")
                text_content = f"[DOCX Document: {uploaded_file.name} - {len(file_content)} bytes - Text extraction failed: {docx_error}]"
                
        elif uploaded_file.name.endswith('.pptx'):
            # Enhanced PPTX processing
            try:
                from pptx import Presentation
                import tempfile
                import os
                
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                    tmp_file.write(file_content)
                    tmp_path = tmp_file.name
                
                try:
                    prs = Presentation(tmp_path)
                    slides_text = []
                    
                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_content = f"\n--- Slide {slide_num} ---\n"
                        
                        # Extract text from shapes
                        text_found = False
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_content += f"{shape.text.strip()}\n"
                                text_found = True
                            
                            # Extract text from tables
                            if hasattr(shape, "table"):
                                table = shape.table
                                for row in table.rows:
                                    row_text = []
                                    for cell in row.cells:
                                        if cell.text.strip():
                                            row_text.append(cell.text.strip())
                                    if row_text:
                                        slide_content += " | ".join(row_text) + "\n"
                                        text_found = True
                        
                        # Extract notes
                        if slide.notes_slide and slide.notes_slide.notes_text_frame:
                            notes_text = slide.notes_slide.notes_text_frame.text.strip()
                            if notes_text:
                                slide_content += f"\nNotes: {notes_text}\n"
                                text_found = True
                        
                        if text_found:
                            slides_text.append(slide_content)
                    
                    text_content = "\n".join(slides_text)
                    
                    if not text_content.strip():
                        text_content = f"[PPTX Document: {uploaded_file.name} - {len(prs.slides)} slides - No extractable text found]"
                        
                finally:
                    os.unlink(tmp_path)
                    
            except Exception as pptx_error:
                logger.warning(f"PPTX processing failed: {pptx_error}")
                text_content = f"[PPTX Document: {uploaded_file.name} - {len(file_content)} bytes - Text extraction failed: {pptx_error}]"
                
        elif (uploaded_file.type and uploaded_file.type.startswith('image/')) or uploaded_file.name.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp')):
            # Enhanced image processing with OCR
            try:
                from PIL import Image
                import tempfile
                import os
                
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                    tmp_file.write(file_content)
                    tmp_path = tmp_file.name
                
                try:
                    # Open and process image
                    image = Image.open(tmp_path)
                    
                    # Get image info
                    image_info = f"Image: {uploaded_file.name}\nSize: {image.size[0]}x{image.size[1]} pixels\nFormat: {image.format}\n\n"
                    
                    # Try OCR text extraction
                    try:
                        import pytesseract
                        
                        # Convert to RGB if necessary
                        if image.mode != 'RGB':
                            image = image.convert('RGB')
                        
                        # Extract text using OCR
                        extracted_text = pytesseract.image_to_string(image, config='--psm 6')
                        
                        if extracted_text.strip():
                            text_content = f"{image_info}OCR Extracted Text:\n{extracted_text.strip()}"
                        else:
                            text_content = f"{image_info}[No text detected in image]"
                            
                    except ImportError:
                        text_content = f"{image_info}[OCR not available - install pytesseract for text extraction]"
                    except Exception as ocr_error:
                        text_content = f"{image_info}[OCR failed: {ocr_error}]"
                        
                finally:
                    os.unlink(tmp_path)
                    
            except Exception as image_error:
                logger.warning(f"Image processing failed: {image_error}")
                text_content = f"[Image: {uploaded_file.name} - {len(file_content)} bytes - Processing failed: {image_error}]"
                
        else:
            text_content = f"[Document: {uploaded_file.name} - {len(file_content)} bytes - Unsupported format. Supported: TXT, MD, PDF, DOCX, PPTX, Images]"
        
        return text_content
        
    except Exception as e:
        logger.error(f"Error processing document: {e}")
        return f"[Error processing {uploaded_file.name}: {e}]"

def save_document_to_conversation(uploaded_file, content):
    """Save document to conversation knowledge base with advanced RAG processing."""
    try:
        if not st.session_state.get('current_conversation_id'):
            return False
        
        conv_id = st.session_state.current_conversation_id
        
        # Save to session state for immediate access
        if 'conversation_documents' not in st.session_state:
            st.session_state.conversation_documents = {}
        
        if conv_id not in st.session_state.conversation_documents:
            st.session_state.conversation_documents[conv_id] = []
        
        # Add document to conversation knowledge base
        document_info = {
            'filename': uploaded_file.name,
            'content': content,
            'uploaded_at': datetime.now().isoformat(),
            'file_size': len(uploaded_file.getvalue())
        }
        
        st.session_state.conversation_documents[conv_id].append(document_info)
        
        # Process with advanced RAG for full document knowledge base
        try:
            from services.rag_service import RAGService
            from services.user_service import user_service
            import uuid
            
            # Get user UUID
            user_data = run_async(user_service.get_user_by_id(st.session_state.user_id))
            if user_data:
                rag_service = RAGService()
                document_id = str(uuid.uuid4())
                
                # Always process document for full knowledge base (default behavior)
                success = run_async(rag_service.process_document(
                    content, 
                    document_id, 
                    conv_id, 
                    user_data['id'],
                    metadata={
                        'filename': uploaded_file.name,
                        'file_size': len(uploaded_file.getvalue()),
                        'uploaded_at': datetime.now().isoformat(),
                        'processing_mode': 'full_document_knowledge_base_default'
                    },
                    use_full_document_mode=True  # Explicitly set to True as default
                ))
                
                if success:
                    logger.info(f"✅ Document processed for full knowledge base: {uploaded_file.name}")
                else:
                    logger.warning(f"⚠️ Full document processing failed for: {uploaded_file.name}")
            
        except Exception as rag_error:
            logger.warning(f"Full document RAG processing failed: {rag_error}")
            # Continue with basic storage
        
        logger.info(f"Document saved to conversation {conv_id}: {uploaded_file.name}")
        return True
        
    except Exception as e:
        logger.error(f"Error saving document to conversation: {e}")
        return False

def get_conversation_context(user_query=""):
    """Get full document context for current conversation (entire documents as knowledge base)."""
    try:
        conv_id = st.session_state.get('current_conversation_id')
        if not conv_id:
            return ""
        
        # Always use full document RAG as primary method (default behavior)
        try:
            from services.rag_service import RAGService
            from services.user_service import user_service
            
            # Get user UUID
            user_data = run_async(user_service.get_user_by_id(st.session_state.user_id))
            if user_data:
                rag_service = RAGService()
                
                # ALWAYS try full document context first (default behavior)
                full_context = run_async(rag_service.get_full_document_context(
                    conv_id, user_data['id'], max_context_length=20000  # Increased for more complete context
                ))
                
                if full_context:
                    logger.info(f"✅ Retrieved complete document knowledge base: {len(full_context)} chars")
                    return f"\n\n--- COMPLETE DOCUMENT KNOWLEDGE BASE (DEFAULT) ---\n{full_context}\n"
                
                # Only use similarity search as last resort
                if user_query:
                    logger.info("Full document context not available, trying similarity search as fallback")
                    similarity_context = run_async(rag_service.get_conversation_context(
                        user_query, conv_id, user_data['id'], max_context_length=10000
                    ))
                    
                    if similarity_context:
                        logger.info(f"✅ Retrieved similarity-based fallback context: {len(similarity_context)} chars")
                        return f"\n\n--- DOCUMENT EXCERPTS (FALLBACK) ---\n{similarity_context}\n"
                    
        except Exception as rag_error:
            logger.warning(f"Full document RAG failed, falling back to session context: {rag_error}")
        
        # Fallback to session-based documents (enhanced)
        if 'conversation_documents' in st.session_state:
            documents = st.session_state.conversation_documents.get(conv_id, [])
            if documents:
                context = "\n\n--- CONVERSATION DOCUMENTS ---\n"
                for doc in documents:
                    # Use more content and better formatting
                    content = doc['content']
                    if len(content) > 2000:
                        # For long documents, try to find relevant sections
                        if user_query:
                            relevant_section = find_relevant_section(content, user_query)
                            context += f"\n[Document: {doc['filename']}]\n{relevant_section}\n"
                        else:
                            # Show beginning and end for context
                            context += f"\n[Document: {doc['filename']}]\n{content[:1500]}...\n[End excerpt: ...{content[-500:]}]\n"
                    else:
                        context += f"\n[Document: {doc['filename']}]\n{content}\n"
                
                return context
        
        return ""
        
    except Exception as e:
        logger.error(f"Error getting conversation context: {e}")
        return ""

def extract_key_terms(query, max_terms=3):
    """Extract key terms from user query for multi-query RAG."""
    try:
        import re
        
        # Remove common stop words
        stop_words = {
            'what', 'how', 'why', 'when', 'where', 'who', 'which', 'is', 'are', 'was', 'were',
            'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with',
            'by', 'from', 'up', 'about', 'into', 'through', 'during', 'before', 'after',
            'above', 'below', 'can', 'could', 'should', 'would', 'will', 'shall', 'may', 'might',
            'do', 'does', 'did', 'have', 'has', 'had', 'be', 'been', 'being', 'this', 'that',
            'these', 'those', 'i', 'you', 'he', 'she', 'it', 'we', 'they', 'me', 'him', 'her',
            'us', 'them', 'my', 'your', 'his', 'her', 'its', 'our', 'their'
        }
        
        # Extract words (3+ characters, not stop words)
        words = re.findall(r'\b[a-zA-Z]{3,}\b', query.lower())
        key_terms = [word for word in words if word not in stop_words]
        
        # Return unique terms, prioritizing longer ones
        unique_terms = list(dict.fromkeys(key_terms))  # Preserve order, remove duplicates
        unique_terms.sort(key=len, reverse=True)  # Longer terms first
        
        return unique_terms[:max_terms]
        
    except Exception as e:
        logger.error(f"Error extracting key terms: {e}")
        return []

def find_relevant_section(content, query, section_size=1500):
    """Find the most relevant section of a document based on query."""
    try:
        query_words = query.lower().split()
        content_lower = content.lower()
        
        # Find positions of query words
        word_positions = []
        for word in query_words:
            pos = content_lower.find(word)
            if pos != -1:
                word_positions.append(pos)
        
        if not word_positions:
            # No matches found, return beginning
            return content[:section_size] + "..." if len(content) > section_size else content
        
        # Find the center of query word cluster
        center_pos = sum(word_positions) // len(word_positions)
        
        # Extract section around the center
        start_pos = max(0, center_pos - section_size // 2)
        end_pos = min(len(content), start_pos + section_size)
        
        # Adjust to word boundaries
        if start_pos > 0:
            # Find previous sentence or paragraph break
            for i in range(start_pos, max(0, start_pos - 100), -1):
                if content[i] in '.!?\n':
                    start_pos = i + 1
                    break
        
        if end_pos < len(content):
            # Find next sentence or paragraph break
            for i in range(end_pos, min(len(content), end_pos + 100)):
                if content[i] in '.!?\n':
                    end_pos = i + 1
                    break
        
        section = content[start_pos:end_pos].strip()
        
        # Add context indicators
        prefix = "..." if start_pos > 0 else ""
        suffix = "..." if end_pos < len(content) else ""
        
        return f"{prefix}{section}{suffix}"
        
    except Exception as e:
        logger.error(f"Error finding relevant section: {e}")
        return content[:section_size] + "..." if len(content) > section_size else content

def render_simple_chatbot():
    """Render a simple, clean chatbot interface."""
    
    # Check authentication
    if not st.session_state.get('authenticated'):
        st.error("Please sign in to use the chatbot.")
        return
    
    st.title("💊 PharmGPT")
    
    # Initialize conversation ID
    if 'current_conversation_id' not in st.session_state:
        st.session_state.current_conversation_id = None
    
    # Initialize model preference
    if 'selected_model_mode' not in st.session_state:
        st.session_state.selected_model_mode = "normal"
    
    # Track conversation changes and load appropriate messages
    if 'last_loaded_conversation' not in st.session_state:
        st.session_state.last_loaded_conversation = None
    
    # Check if conversation changed
    current_conv_id = st.session_state.get('current_conversation_id')
    if st.session_state.last_loaded_conversation != current_conv_id:
        load_conversation_messages()
        st.session_state.last_loaded_conversation = current_conv_id
        # Clear any pending document uploads when switching conversations
        if 'pending_document' in st.session_state:
            del st.session_state.pending_document
    
    # Initialize messages if not already done
    if 'chat_messages' not in st.session_state:
        st.session_state.chat_messages = []
    
    # Simple model selection
    col1, col2, col3 = st.columns([2, 1, 2])
    with col2:
        is_turbo = st.toggle("⚡ Turbo Mode", 
                           value=(st.session_state.selected_model_mode == "turbo"),
                           help="Switch between Normal and Turbo modes")
        st.session_state.selected_model_mode = "turbo" if is_turbo else "normal"
    
    # Show enhanced document status for current conversation
    conv_id = st.session_state.get('current_conversation_id')
    if conv_id and 'conversation_documents' in st.session_state:
        documents = st.session_state.conversation_documents.get(conv_id, [])
        if documents:
            doc_count = len(documents)
            total_chars = sum(len(doc.get('content', '')) for doc in documents)
            
            # Show document summary
            doc_names = [doc.get('filename', 'Unknown') for doc in documents]
            doc_list = ", ".join(doc_names[:3])  # Show first 3 document names
            if len(doc_names) > 3:
                doc_list += f" and {len(doc_names) - 3} more"
            
            # Documents loaded - no verbose status needed
    
    # Display chat messages with regenerate button
    for i, message in enumerate(st.session_state.chat_messages):
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
            
            # Add regenerate button for the last assistant message
            if (message["role"] == "assistant" and 
                i == len(st.session_state.chat_messages) - 1 and 
                i > 0):  # Only show for the last assistant message
                
                if st.button("🔄 Regenerate", key=f"regen_{i}", help="Regenerate response"):
                    # Get the previous user message
                    if i > 0 and st.session_state.chat_messages[i-1]["role"] == "user":
                        user_prompt = st.session_state.chat_messages[i-1]["content"]
                        
                        # Remove the current assistant response
                        st.session_state.chat_messages.pop()
                        
                        # Set regeneration flag
                        st.session_state.regenerate_response = True
                        st.session_state.regenerate_prompt = user_prompt
                        st.rerun()
    
    # Document upload and chat input area
    st.markdown("---")
    
    # Document upload beside message input
    col1, col2 = st.columns([3, 1])
    
    with col1:
        # Chat input
        prompt = st.chat_input("Ask me anything about pharmacology...")
    
    with col2:
        # Document upload with enhanced format support
        uploaded_file = st.file_uploader(
            "📎 Upload",
            type=['txt', 'pdf', 'docx', 'md', 'pptx', 'png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'webp'],
            help="Upload documents, presentations, or images for this conversation",
            key=f"doc_upload_{st.session_state.get('current_conversation_id', 'new')}"
        )
    
    # Handle document upload
    if uploaded_file is not None:
        with st.spinner(f"Processing {uploaded_file.name}..."):
            document_content = process_uploaded_document(uploaded_file)
            if document_content:
                success = save_document_to_conversation(uploaded_file, document_content)
                if success:
                    st.rerun()
                else:
                    st.error("❌ Failed to save document")
            else:
                st.error("❌ Failed to process document")
    
    # Handle regenerate response
    if st.session_state.get('regenerate_response', False):
        prompt = st.session_state.get('regenerate_prompt', '')
        st.session_state.regenerate_response = False
        st.session_state.regenerate_prompt = None
    
    # Process chat input or regeneration
    if prompt:
        logger.info(f"User input: {prompt[:50]}...")
        
        # Add user message to chat history
        st.session_state.chat_messages.append({"role": "user", "content": prompt})
        
        # Display user message
        with st.chat_message("user"):
            st.markdown(prompt)
        
        # Generate assistant response
        with st.chat_message("assistant"):
            response_placeholder = st.empty()
            full_response = ""
            
            try:
                # Import API functions
                from openai_client import chat_completion_stream, chat_completion, get_available_model_modes
                from prompts import pharmacology_system_prompt
                
                # Get available models
                available_modes = get_available_model_modes()
                if not available_modes:
                    response_placeholder.markdown("❌ No API models available. Please check your API keys.")
                    return
                
                # Use selected model mode
                selected_mode = st.session_state.selected_model_mode
                if selected_mode not in available_modes:
                    selected_mode = list(available_modes.keys())[0]
                
                model = available_modes[selected_mode]["model"]
                logger.info(f"Using {selected_mode} mode: {model}")
                
                # Get enhanced document context for this conversation
                document_context = get_conversation_context(prompt)
                
                # Prepare system prompt with complete document knowledge base (default behavior)
                system_prompt = pharmacology_system_prompt
                if document_context:
                    system_prompt += f"\n\n{document_context}\n\nThe above contains the complete content of all documents uploaded to this conversation. This is your primary knowledge source. Always search the document content first before using general knowledge, quote specific sections when referencing information, and prioritize document information over your training data when conflicts arise."
                
                # Prepare messages for API
                api_messages = [{"role": "system", "content": system_prompt}]
                
                # Add recent conversation history (last 10 messages)
                for msg in st.session_state.chat_messages[-10:]:
                    api_messages.append({"role": msg["role"], "content": msg["content"]})
                
                # Ultra-fluid streaming with 7500 token capacity
                try:
                    response_placeholder.markdown("💭 Thinking...")
                    
                    for chunk in chat_completion_stream(model, api_messages):
                        if chunk:
                            full_response += chunk
                            # Ultra-fluid: update on every chunk for maximum responsiveness
                            response_placeholder.markdown(full_response + "▌")
                    
                    # Clean final display
                    if full_response.strip():
                        response_placeholder.markdown(full_response)
                        logger.info(f"✅ Streaming completed: {len(full_response)} chars")
                    else:
                        raise Exception("Streaming failed or empty response")
                        
                except Exception as stream_error:
                    logger.warning(f"Streaming failed: {stream_error}, trying fallback...")
                    
                    # Fallback to non-streaming
                    response_placeholder.markdown("💭 Processing...")
                    full_response = chat_completion(model, api_messages)
                    
                    if full_response and not full_response.startswith("Error:"):
                        response_placeholder.markdown(full_response)
                        logger.info(f"✅ Fallback completed: {len(full_response)} chars")
                    else:
                        raise Exception(f"API call failed: {full_response}")
                
                # Add response to chat history
                if full_response and not full_response.startswith("Error:") and not full_response.startswith("❌"):
                    st.session_state.chat_messages.append({"role": "assistant", "content": full_response})
                    logger.info(f"✅ Response added to chat history ({len(full_response)} chars)")
                    
                    # Auto-save conversation
                    try:
                        save_success = save_conversation()
                        if save_success:
                            logger.info("✅ Conversation auto-saved")
                    except Exception as save_error:
                        logger.error(f"Auto-save error: {save_error}")
                else:
                    logger.error(f"Response not added: {full_response[:100] if full_response else 'Empty response'}...")
                    
            except Exception as e:
                error_msg = f"❌ Error generating response: {str(e)}"
                response_placeholder.markdown(error_msg)
                logger.error(f"Exception in response generation: {e}")

def render_chatbot_page():
    """Main function called by the app."""
    render_simple_chatbot()