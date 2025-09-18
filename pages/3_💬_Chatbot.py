"""
Chatbot Page - Multipage Streamlit App with Fixed Message Input Position
"""

import streamlit as st
import sys
import os
import logging
from datetime import datetime
import functools
from pathlib import Path

# Get parent directory path
current_dir = Path(__file__).parent
parent_dir = current_dir.parent

# Add to Python path
sys.path.insert(0, str(parent_dir))

# Import required modules
from utils.session_manager import initialize_session_state
from utils.theme import apply_theme
from auth import initialize_auth_session, logout_current_user
from config import APP_TITLE, APP_ICON, MAX_FILE_SIZE_MB

# Configure logging
logger = logging.getLogger(__name__)

# Mobile-responsive page configuration
st.set_page_config(
    page_title=f"{APP_TITLE} - Chat",
    page_icon=APP_ICON,
    layout="wide",
    initial_sidebar_state="auto"  # Auto-collapse on mobile
)

def main():
    """Chatbot page entry point with optimized performance."""
    # Initialize session state and authentication
    initialize_session_state()
    apply_theme()
    initialize_auth_session()
    
    # Add mobile-specific enhancements and theme-aware styling
    st.markdown("""
    <style>
        /* Enhanced toggle styling for model selection */
        .stToggle > div > label {
            font-weight: 500 !important;
            font-size: 16px !important;
        }
        
        /* Theme-aware welcome card styling */
        .welcome-card {
            transition: all 0.3s ease !important;
        }
        
        .welcome-tag {
            transition: all 0.3s ease !important;
        }
        
        /* Smooth scrolling for streaming responses */
        html {
            scroll-behavior: smooth !important;
        }
        
        .main {
            scroll-behavior: smooth !important;
        }
        
        /* Ensure streaming responses are visible */
        .stMarkdown {
            scroll-margin-bottom: 100px !important;
        }
        
        /* Auto-scroll container */
        .streaming-container {
            scroll-margin-top: 20px !important;
            scroll-margin-bottom: 20px !important;
        }
        
        /* Horizontal scrollable tables */
        .stMarkdown table {
            display: block !important;
            overflow-x: auto !important;
            white-space: nowrap !important;
            max-width: 100% !important;
            border-collapse: collapse !important;
            margin: 1rem 0 !important;
        }
        
        .stMarkdown table thead,
        .stMarkdown table tbody,
        .stMarkdown table tr {
            display: table !important;
            width: 100% !important;
            table-layout: fixed !important;
        }
        
        .stMarkdown table th,
        .stMarkdown table td {
            padding: 8px 12px !important;
            border: 1px solid #e5e7eb !important;
            text-align: left !important;
            min-width: 120px !important;
            white-space: nowrap !important;
            overflow: hidden !important;
            text-overflow: ellipsis !important;
        }
        
        .stMarkdown table th {
            background-color: #f8fafc !important;
            font-weight: 600 !important;
            position: sticky !important;
            top: 0 !important;
        }
        
        /* Dark mode table styling */
        @media (prefers-color-scheme: dark) {
            .stMarkdown table th,
            .stMarkdown table td {
                border-color: #475569 !important;
            }
            
            .stMarkdown table th {
                background-color: #1e293b !important;
                color: #f8fafc !important;
            }
            
            .stMarkdown table td {
                color: #f8fafc !important;
            }
        }
        
        /* Table container with scroll indicators */
        .stMarkdown table {
            box-shadow: 0 0 0 1px #e5e7eb !important;
            border-radius: 8px !important;
        }
        
        /* Scroll hint for tables */
        .stMarkdown table::after {
            content: "← Scroll horizontally to see more →" !important;
            display: block !important;
            text-align: center !important;
            font-size: 12px !important;
            color: #6b7280 !important;
            padding: 8px !important;
            background-color: #f9fafb !important;
            border-top: 1px solid #e5e7eb !important;
        }
        
        /* Mobile chat optimizations */
        @media (max-width: 768px) {
            /* Hide sidebar toggle on mobile for cleaner look */
            .stSidebarNav {
                display: none !important;
            }
            
            /* Optimize chat input for mobile - seamless blend */
            .stChatInputContainer {
                position: fixed !important;
                bottom: 0 !important;
                left: 0 !important;
                right: 0 !important;
                background: transparent !important;
                border: none !important;
                padding: 0.5rem !important;
                z-index: 1000 !important;
            }
            
            /* Add bottom padding to prevent content hiding behind fixed input */
            .main .block-container {
                padding-bottom: 80px !important;
            }
            
            /* Mobile-friendly message display */
            .stChatMessage {
                word-wrap: break-word !important;
                overflow-wrap: break-word !important;
            }
            
            /* Mobile welcome card adjustments */
            .welcome-card {
                padding: 1.5rem !important;
                margin: 0.5rem 0 !important;
            }
            
            .welcome-tag {
                padding: 0.4rem 0.8rem !important;
                font-size: 14px !important;
            }
        }
        
        /* Touch-friendly improvements */
        .stButton > button:active {
            transform: scale(0.98) !important;
            transition: transform 0.1s !important;
        }
        
        /* Improve readability on small screens */
        @media (max-width: 480px) {
            .stMarkdown {
                font-size: 14px !important;
                line-height: 1.6 !important;
            }
            
            h1, h2, h3 {
                font-size: 1.2em !important;
                margin-bottom: 0.5rem !important;
            }
        }
    </style>
    """, unsafe_allow_html=True)
    
    # Check authentication first (faster than validation)
    if not st.session_state.get('authenticated'):
        st.error("🔐 Please sign in to access the chatbot.")
        col1, col2, col3 = st.columns([1, 2, 1])
        with col2:
            if st.button("🔐 Go to Sign In", use_container_width=True, type="primary"):
                st.switch_page("pages/2_🔐_Sign_In.py")
            if st.button("← Back to Home", use_container_width=True):
                st.switch_page("app.py")
        return
    
    # Basic session validation without aggressive clearing
    if not st.session_state.get('authenticated'):
        st.error("🔐 Please sign in to access the chatbot.")
        col1, col2, col3 = st.columns([1, 2, 1])
        with col2:
            if st.button("🔐 Go to Sign In", use_container_width=True, type="primary"):
                st.switch_page("pages/2_🔐_Sign_In.py")
            if st.button("← Back to Home", use_container_width=True):
                st.switch_page("app.py")
        return
    
    # Render the chatbot interface
    render_chatbot_interface()

def render_chatbot_interface():
    """Render the enhanced chatbot interface with optimized performance."""
    
    # Ensure conversations are loaded and user_id is correct
    ensure_conversations_loaded()
    
    # Enhanced sidebar for conversation management
    render_enhanced_sidebar()
    
    # Main chat area with improved layout
    render_main_chat_area()

def ensure_conversations_loaded():
    """Ensure conversations are properly loaded with correct user_id."""
    try:
        # Check if user_id is set correctly
        username = st.session_state.get('username')
        user_id = st.session_state.get('user_id')
        
        if username and not user_id:
            logger.warning(f"Missing user_id for username: {username}, attempting to fix...")
            from auth import get_user_legacy_id
            correct_user_id = get_user_legacy_id(username)
            if correct_user_id:
                st.session_state.user_id = correct_user_id
                logger.info(f"Fixed user_id to: {correct_user_id}")
        
        # Ensure conversations are loaded
        from fix_user_isolation import get_secure_conversations
        conversations = get_secure_conversations()
        
        logger.info(f"Conversations loaded: {len(conversations)}")
        
    except Exception as e:
        logger.error(f"Error ensuring conversations loaded: {e}")

def load_user_conversations():
    """Load user conversations from database with proper user isolation and caching."""
    if not st.session_state.get('authenticated') or not st.session_state.get('user_id'):
        return
    
    try:
        # Load conversations safely without aggressive validation
        from fix_user_isolation import load_user_conversations_safely, secure_update_conversations
        
        # Load conversations safely
        conversations = load_user_conversations_safely()
        secure_update_conversations(conversations)
        
        logger.info(f"Safely loaded {len(conversations)} conversations for user: {st.session_state.username}")
        
        # Log conversation IDs for debugging (first 8 chars only)
        conv_ids = [conv_id[:8] + "..." for conv_id in conversations.keys()]
        logger.info(f"Conversation IDs: {conv_ids}")
        
    except Exception as e:
        logger.error(f"Failed to load conversations: {e}")
        # Don't clear conversations on error, just log it
        logger.warning("Keeping existing conversations due to load error")

def render_enhanced_sidebar():
    """Render enhanced sidebar with conversation management."""
    with st.sidebar:
        # Theme toggle at the top for easy access
        from utils.theme import render_theme_toggle, add_mobile_meta_tags
        add_mobile_meta_tags()
        render_theme_toggle()
        
        # Performance settings right after theme toggle
        st.markdown("### ⚙️ Performance Settings")
        
        # Model selection - Toggle between fast and premium
        use_premium_model = st.toggle(
            "💎 Premium Mode", 
            value=st.session_state.get('selected_model_mode', 'fast') == 'premium',
            help="Toggle between ⚡ Fast Mode (default) and 💎 Premium Mode for higher quality responses"
        )
        st.session_state.selected_model_mode = "premium" if use_premium_model else "fast"
        
        # Set streaming as default (always enabled, 6 tokens per second)
        st.session_state.use_streaming = True
        
        st.markdown("---")
        
        st.markdown("### 💊 PharmGPT")
        st.markdown(f"**Welcome, {st.session_state.username}!**")
        
        # Navigation buttons
        col1, col2 = st.columns(2)
        with col1:
            if st.button("🏠 Home", use_container_width=True):
                st.switch_page("app.py")
        with col2:
            if st.button("🚪 Logout", use_container_width=True):
                logout_current_user()
                st.switch_page("app.py")
        
        # Contact Support button
        if st.button("📞 Contact Support", use_container_width=True):
            st.switch_page("pages/4_📞_Contact_Support.py")
        
        # Admin Panel access (only for admin user)
        if st.session_state.get('username') == 'admin':
            if st.button("🛠️ Admin Panel", use_container_width=True):
                # Set admin mode flag and redirect to contact support page
                st.session_state.admin_mode = True
                st.switch_page("pages/4_📞_Contact_Support.py")
        
        st.markdown("---")
        
        # New conversation button
        if st.button("➕ New Conversation", use_container_width=True, type="primary"):
            create_new_conversation()
            # Force refresh the conversation list
            from fix_user_isolation import load_user_conversations_safely, secure_update_conversations
            fresh_conversations = load_user_conversations_safely()
            if fresh_conversations:
                secure_update_conversations(fresh_conversations)
        

        
        # Conversation list
        render_conversation_list()
        
        # Current conversation info
        render_conversation_info()
        


def create_new_conversation():
    """Create a new conversation and save it to database immediately."""
    try:
        # Create conversation in database immediately
        from utils.conversation_manager import run_async, create_new_conversation as create_db_conversation
        
        conversation_id = run_async(create_db_conversation())
        if conversation_id:
            st.session_state.current_conversation_id = conversation_id
            st.session_state.chat_messages = []
            if 'conversation_documents' in st.session_state:
                st.session_state.conversation_documents = {}
            
            logger.info(f"Created new conversation: {conversation_id}")
            st.success("✅ New conversation created!")
        else:
            # Fallback to old behavior
            st.session_state.chat_messages = []
            st.session_state.current_conversation_id = None
            if 'conversation_documents' in st.session_state:
                st.session_state.conversation_documents = {}
            logger.warning("Failed to create conversation in database, using fallback")
            
    except Exception as e:
        logger.error(f"Error creating new conversation: {e}")
        # Fallback to old behavior
        st.session_state.chat_messages = []
        st.session_state.current_conversation_id = None
        if 'conversation_documents' in st.session_state:
            st.session_state.conversation_documents = {}
    
    st.rerun()

def render_conversation_list():
    """Render the list of conversations."""
    st.markdown("### 💬 Conversations")
    
    from fix_user_isolation import get_secure_conversations
    conversations = get_secure_conversations()
    
    if not conversations:
        st.info("No conversations yet. Start chatting to create your first conversation!")
        return
    
    # Sort conversations by updated_at (most recent first)
    sorted_conversations = sorted(
        conversations.items(),
        key=lambda x: x[1].get('updated_at', x[1].get('created_at', '')),
        reverse=True
    )
    
    # Show recent conversations (limit to 10)
    for conv_id, conv_data in sorted_conversations[:10]:
        title = conv_data.get('title', 'Untitled Chat')
        message_count = len(conv_data.get('messages', []))
        
        # Truncate long titles
        display_title = title[:30] + "..." if len(title) > 30 else title
        
        # Highlight current conversation
        is_current = st.session_state.get('current_conversation_id') == conv_id
        button_type = "primary" if is_current else "secondary"
        
        # Create conversation item with action buttons
        with st.container():
            # Main conversation button
            col1, col2, col3 = st.columns([6, 1, 1])
            
            with col1:
                if st.button(
                    f"💬 {display_title}\n({message_count} messages)",
                    key=f"conv_{conv_id}",
                    use_container_width=True,
                    type=button_type,
                    help=f"Switch to: {title}"
                ):
                    load_conversation(conv_id)
            
            with col2:
                # Delete button for each conversation
                if st.button("🗑️", key=f"delete_{conv_id}", help="Delete conversation"):
                    delete_specific_conversation(conv_id)
            
            # Removed export functionality
    
    # Show more conversations option
    if len(conversations) > 10:
        with st.expander(f"📁 Show {len(conversations) - 10} more conversations"):
            for conv_id, conv_data in sorted_conversations[10:]:
                title = conv_data.get('title', 'Untitled Chat')
                message_count = len(conv_data.get('messages', []))
                display_title = title[:20] + "..." if len(title) > 20 else title
                
                # More conversations with delete buttons
                col1, col2, col3 = st.columns([6, 1, 1])
                
                with col1:
                    if st.button(
                        f"{display_title} ({message_count})",
                        key=f"conv_more_{conv_id}",
                        use_container_width=True
                    ):
                        load_conversation(conv_id)
                
                with col2:
                    if st.button("🗑️", key=f"delete_more_{conv_id}", help="Delete conversation"):
                        delete_specific_conversation(conv_id)
                
                # Removed export functionality

def load_conversation(conversation_id):
    """Load a specific conversation."""
    from fix_user_isolation import get_secure_conversations
    conversations = get_secure_conversations()
    if conversation_id in conversations:
        st.session_state.current_conversation_id = conversation_id
        st.session_state.chat_messages = conversations[conversation_id].get('messages', [])
        logger.info(f"Loaded conversation: {conversation_id}")
        st.rerun()

def render_conversation_info():
    """Render current conversation information."""
    if not st.session_state.get('current_conversation_id'):
        return
    
    st.markdown("---")
    st.markdown("### 📊 Current Chat")
    
    # Message count
    message_count = len(st.session_state.get('chat_messages', []))
    st.metric("Messages", message_count)
    
    # Document count
    conv_id = st.session_state.current_conversation_id
    doc_count = 0
    if 'conversation_documents' in st.session_state:
        doc_count = len(st.session_state.conversation_documents.get(conv_id, []))
    
    if doc_count > 0:
        st.metric("Documents", doc_count)
    
    # Show rename dialog if requested
    if st.session_state.get('show_rename_dialog', False):
        render_rename_dialog()

def delete_specific_conversation(conv_id: str):
    """Delete a specific conversation from the sidebar."""
    try:
        # Confirm deletion with user
        if f"confirm_delete_{conv_id}" not in st.session_state:
            st.session_state[f"confirm_delete_{conv_id}"] = False
        
        if not st.session_state[f"confirm_delete_{conv_id}"]:
            # First click - ask for confirmation
            st.session_state[f"confirm_delete_{conv_id}"] = True
            st.warning("⚠️ Click delete again to confirm")
            return
        
        # Second click - actually delete
        from utils.conversation_manager import run_async, delete_conversation
        success = run_async(delete_conversation(conv_id))
        
        if success:
            # Remove from session state securely
            from fix_user_isolation import secure_delete_conversation, get_secure_conversations
            secure_delete_conversation(conv_id)
            
            # Clean up confirmation state
            del st.session_state[f"confirm_delete_{conv_id}"]
            
            # If this was the current conversation, switch to another or clear
            if st.session_state.get('current_conversation_id') == conv_id:
                remaining_conversations = list(get_secure_conversations().keys())
                if remaining_conversations:
                    # Switch to the most recent conversation
                    st.session_state.current_conversation_id = remaining_conversations[0]
                    conversations = get_secure_conversations()
                    st.session_state.chat_messages = conversations[remaining_conversations[0]].get('messages', [])
                else:
                    # No conversations left
                    st.session_state.current_conversation_id = None
                    st.session_state.chat_messages = []
            
            st.success("✅ Conversation deleted!")
            st.rerun()
        else:
            st.error("❌ Failed to delete conversation")
            st.session_state[f"confirm_delete_{conv_id}"] = False
            
    except Exception as e:
        logger.error(f"Error deleting conversation: {e}")
        st.error(f"❌ Error: {e}")
        if f"confirm_delete_{conv_id}" in st.session_state:
            st.session_state[f"confirm_delete_{conv_id}"] = False

def duplicate_current_conversation():
    """Duplicate the current conversation."""
    conv_id = st.session_state.get('current_conversation_id')
    if not conv_id:
        return
    
    try:
        from utils.conversation_manager import run_async, duplicate_conversation
        new_conv_id = run_async(duplicate_conversation(conv_id))
        
        if new_conv_id:
            st.success("✅ Conversation duplicated!")
            # Reload conversations to show the new one
            load_user_conversations()
            st.rerun()
        else:
            st.error("❌ Failed to duplicate conversation")
    except Exception as e:
        logger.error(f"Error duplicating conversation: {e}")
        st.error(f"❌ Error: {e}")

# Export functions removed

def render_rename_dialog():
    """Render dialog to rename current conversation."""
    conv_id = st.session_state.get('current_conversation_id')
    if not conv_id:
        st.session_state.show_rename_dialog = False
        return
    
    from fix_user_isolation import get_secure_conversations
    conversations = get_secure_conversations()
    current_title = conversations.get(conv_id, {}).get('title', 'Untitled Chat')
    
    st.markdown("**✏️ Rename Conversation**")
    
    new_title = st.text_input(
        "New title:",
        value=current_title,
        key="rename_input",
        placeholder="Enter new conversation title"
    )
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.button("✅ Save", use_container_width=True):
            if new_title and new_title.strip():
                try:
                    from utils.conversation_manager import run_async, update_conversation_title
                    success = run_async(update_conversation_title(conv_id, new_title.strip()))
                    
                    if success:
                        # Update local state securely
                        from fix_user_isolation import secure_update_conversation
                        secure_update_conversation(conv_id, {'title': new_title.strip()})
                        st.success("✅ Conversation renamed!")
                        st.session_state.show_rename_dialog = False
                        st.rerun()
                    else:
                        st.error("❌ Failed to rename conversation")
                except Exception as e:
                    logger.error(f"Error renaming conversation: {e}")
                    st.error(f"❌ Error: {e}")
            else:
                st.error("❌ Please enter a valid title")
    
    with col2:
        if st.button("❌ Cancel", use_container_width=True):
            st.session_state.show_rename_dialog = False
            st.rerun()

# Document processing functions

def process_uploaded_document_multipage(uploaded_file):
    """Process uploaded document and return text content with enhanced support for images and PPTX."""
    try:
        # Read file content
        file_content = uploaded_file.read()
        uploaded_file.seek(0)  # Reset file pointer
        
        # Process based on file type
        if uploaded_file.type == "text/plain" or uploaded_file.name.endswith('.txt'):
            text_content = file_content.decode('utf-8')
            
        elif uploaded_file.name.endswith('.md'):
            text_content = file_content.decode('utf-8')
            
        elif uploaded_file.type == "application/pdf" or uploaded_file.name.endswith('.pdf'):
            try:
                import PyPDF2
                import io
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                text_content = ""
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content += f"\n--- Page {page_num} ---\n{page_text}\n"
                
                if not text_content.strip():
                    text_content = f"[PDF Document: {uploaded_file.name} - {len(file_content)} bytes - No extractable text found]"
                    
            except Exception as pdf_error:
                logger.warning(f"PDF processing failed: {pdf_error}")
                text_content = f"[PDF Document: {uploaded_file.name} - {len(file_content)} bytes - Text extraction failed: {pdf_error}]"
                
        elif uploaded_file.name.endswith('.docx'):
            try:
                import docx2txt
                import tempfile
                import os
                
                with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
                    tmp_file.write(file_content)
                    tmp_path = tmp_file.name
                
                try:
                    text_content = docx2txt.process(tmp_path)
                    if not text_content.strip():
                        # Try alternative method with python-docx
                        try:
                            from docx import Document
                            doc = Document(tmp_path)
                            paragraphs = []
                            for paragraph in doc.paragraphs:
                                if paragraph.text.strip():
                                    paragraphs.append(paragraph.text)
                            text_content = "\n".join(paragraphs)
                        except:
                            raise Exception("No text extracted with any method")
                finally:
                    os.unlink(tmp_path)
                    
            except Exception as docx_error:
                logger.warning(f"DOCX processing failed: {docx_error}")
                text_content = f"[DOCX Document: {uploaded_file.name} - {len(file_content)} bytes - Text extraction failed: {docx_error}]"
                
        elif uploaded_file.name.endswith('.pptx'):
            # Enhanced PPTX processing
            try:
                from pptx import Presentation
                import tempfile
                import os
                
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                    tmp_file.write(file_content)
                    tmp_path = tmp_file.name
                
                try:
                    prs = Presentation(tmp_path)
                    slides_text = []
                    
                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_content = f"\n--- Slide {slide_num} ---\n"
                        
                        # Extract text from shapes
                        text_found = False
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_content += f"{shape.text.strip()}\n"
                                text_found = True
                            
                            # Extract text from tables
                            if hasattr(shape, "table"):
                                table = shape.table
                                for row in table.rows:
                                    row_text = []
                                    for cell in row.cells:
                                        if cell.text.strip():
                                            row_text.append(cell.text.strip())
                                    if row_text:
                                        slide_content += " | ".join(row_text) + "\n"
                                        text_found = True
                        
                        # Extract notes
                        if slide.notes_slide and slide.notes_slide.notes_text_frame:
                            notes_text = slide.notes_slide.notes_text_frame.text.strip()
                            if notes_text:
                                slide_content += f"\nNotes: {notes_text}\n"
                                text_found = True
                        
                        if text_found:
                            slides_text.append(slide_content)
                    
                    text_content = "\n".join(slides_text)
                    
                    if not text_content.strip():
                        text_content = f"[PPTX Document: {uploaded_file.name} - {len(prs.slides)} slides - No extractable text found]"
                        
                finally:
                    os.unlink(tmp_path)
                    
            except Exception as pptx_error:
                logger.warning(f"PPTX processing failed: {pptx_error}")
                text_content = f"[PPTX Document: {uploaded_file.name} - {len(file_content)} bytes - Text extraction failed: {pptx_error}]"
                
        elif (uploaded_file.type and uploaded_file.type.startswith('image/')) or uploaded_file.name.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp')):
            # Enhanced image processing with multiple OCR engines
            try:
                from utils.ocr_manager import process_image_file
                text_content = process_image_file(uploaded_file, file_content)
                
            except Exception as image_error:
                logger.warning(f"Image processing failed: {image_error}")
                # Fallback to basic image info
                try:
                    from PIL import Image
                    import tempfile
                    import os
                    
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                        tmp_file.write(file_content)
                        tmp_path = tmp_file.name
                    
                    try:
                        image = Image.open(tmp_path)
                        image_info = f"Image: {uploaded_file.name}\nSize: {image.size[0]}x{image.size[1]} pixels\nFormat: {image.format}\n\n"
                        text_content = f"{image_info}IMAGE UPLOADED: OCR processing failed. Please describe what you see in the image and I can help explain the pharmacological concepts shown."
                    finally:
                        os.unlink(tmp_path)
                        
                except Exception as fallback_error:
                    text_content = f"[Image: {uploaded_file.name} - {len(file_content)} bytes - Processing failed: {image_error}]"
                
        else:
            text_content = f"[Document: {uploaded_file.name} - {len(file_content)} bytes - Unsupported format. Supported: TXT, MD, PDF, DOCX, PPTX, Images]"
        
        return text_content
        
    except Exception as e:
        logger.error(f"Error processing document: {e}")
        return f"[Error processing {uploaded_file.name}: {e}]"

def save_document_to_conversation_multipage(uploaded_file, content):
    """Save document to conversation knowledge base with advanced RAG processing."""
    try:
        logger.info(f"Starting document save for: {uploaded_file.name}")
        
        # Check if we have a conversation ID
        if not st.session_state.get('current_conversation_id'):
            logger.error("No current conversation ID found")
            # Create a new conversation if none exists
            try:
                from utils.conversation_manager import run_async, create_new_conversation
                conv_id = run_async(create_new_conversation())
                if not conv_id:
                    logger.error("Failed to create new conversation for document")
                    return False
                st.session_state.current_conversation_id = conv_id
                logger.info(f"Created new conversation for document: {conv_id}")
            except Exception as conv_error:
                logger.error(f"Error creating conversation: {conv_error}")
                return False
        
        conv_id = st.session_state.current_conversation_id
        logger.info(f"Using conversation ID: {conv_id}")
        
        # Check authentication
        if not st.session_state.get('authenticated') or not st.session_state.get('user_id'):
            logger.error("User not authenticated")
            return False
        
        # Save to session state for immediate access
        if 'conversation_documents' not in st.session_state:
            st.session_state.conversation_documents = {}
        
        if conv_id not in st.session_state.conversation_documents:
            st.session_state.conversation_documents[conv_id] = []
        
        # Add document to conversation knowledge base
        try:
            document_info = {
                'filename': uploaded_file.name,
                'content': content,
                'uploaded_at': datetime.now().isoformat(),
                'file_size': len(uploaded_file.getvalue())
            }
            
            st.session_state.conversation_documents[conv_id].append(document_info)
            logger.info(f"Document added to session state: {uploaded_file.name}")
        except Exception as session_error:
            logger.error(f"Error adding document to session state: {session_error}")
            return False
        
        # Try RAG processing (optional - don't fail if this doesn't work)
        try:
            from services.rag_service import RAGService
            from services.user_service import user_service
            from utils.conversation_manager import run_async
            import uuid
            
            logger.info("Starting RAG processing...")
            
            # Get user UUID
            user_data = run_async(user_service.get_user_by_id(st.session_state.user_id))
            if user_data:
                rag_service = RAGService()
                document_id = str(uuid.uuid4())
                
                # Always process document for full knowledge base (default behavior)
                success = run_async(rag_service.process_document(
                    content, 
                    document_id, 
                    conv_id, 
                    user_data['id'],
                    metadata={
                        'filename': uploaded_file.name,
                        'file_size': len(uploaded_file.getvalue()),
                        'uploaded_at': datetime.now().isoformat(),
                        'processing_mode': 'full_document_knowledge_base_default'
                    },
                    use_full_document_mode=True  # Explicitly set to True as default
                ))
                
                if success:
                    logger.info(f"✅ Document processed for full knowledge base: {uploaded_file.name}")
                else:
                    logger.warning(f"⚠️ Full document processing failed for: {uploaded_file.name}")
            else:
                logger.warning("User data not found for RAG processing")
            
        except Exception as rag_error:
            logger.warning(f"RAG processing failed (non-critical): {rag_error}")
            # Continue - RAG failure shouldn't prevent document saving
        
        logger.info(f"Document successfully saved to conversation {conv_id}: {uploaded_file.name}")
        return True
        
    except Exception as e:
        logger.error(f"Critical error saving document to conversation: {e}")
        import traceback
        logger.error(f"Traceback: {traceback.format_exc()}")
        return False

def render_main_chat_area():
    """Render the main chat area with enhanced features."""
    
    # Chat header with conversation title
    render_chat_header()
    
    # Initialize chat messages
    if 'chat_messages' not in st.session_state:
        st.session_state.chat_messages = []
    
    # Show welcome message for new conversations
    if not st.session_state.chat_messages:
        render_welcome_message()
    
    # Create a container for messages that will scroll
    message_container = st.container()
    
    with message_container:
        # Display chat messages with enhanced features
        render_chat_messages()
        
        # Show processing indicator if generating response
        if st.session_state.get('processing_input', False):
            with st.chat_message("assistant"):
                st.markdown("🤔 Thinking...")
                st.caption("Generating response...")
    
    # Fixed bottom input area
    render_bottom_input_area()
    
    # Force refresh if we just added a message
    if st.session_state.get('force_refresh', False):
        st.session_state.force_refresh = False
        st.rerun()

def render_chat_header():
    """Render the chat header with conversation title."""
    col1, col2, col3 = st.columns([3, 1, 1])
    
    with col1:
        if st.session_state.get('current_conversation_id'):
            from fix_user_isolation import get_secure_conversations
            conversations = get_secure_conversations()
            conv_id = st.session_state.current_conversation_id
            title = conversations.get(conv_id, {}).get('title', 'Untitled Chat')
            st.title(f"💊 {title}")
        else:
            st.title("💊 PharmGPT Chat")
    
    with col2:
        # Model indicator
        mode = st.session_state.get('selected_model_mode', 'normal')
        mode_icon = "⚡" if mode == "turbo" else "🧠"
        st.markdown(f"**{mode_icon} {mode.title()} Mode**")
    
    with col3:
        # Document indicator
        conv_id = st.session_state.get('current_conversation_id')
        if conv_id and 'conversation_documents' in st.session_state:
            doc_count = len(st.session_state.conversation_documents.get(conv_id, []))
            if doc_count > 0:
                st.markdown(f"**📎 {doc_count} docs**")

def render_welcome_message():
    """Render welcome message for new conversations."""
    st.markdown("""
    <div class="welcome-card" style="
        text-align: center; 
        padding: 2rem; 
        background: linear-gradient(135deg, #f0f9ff 0%, #e0f2fe 100%); 
        border-radius: 12px; 
        margin: 1rem 0;
        border: 1px solid #e0e7ff;
        color: #1f2937;
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1);
        transition: all 0.3s ease;
    ">
        <h3 style="color: inherit; margin-bottom: 1rem;">👋 Welcome to PharmGPT!</h3>
        <p style="color: inherit; margin-bottom: 1.5rem;">I'm your AI pharmacology assistant. Ask me anything about:</p>
        <div style="display: flex; justify-content: center; gap: 1rem; flex-wrap: wrap; margin-top: 1rem;">
            <span class="welcome-tag" style="background: #ffffff; padding: 0.5rem 1rem; border-radius: 20px; border: 1px solid #e0e7ff; color: inherit; transition: all 0.3s ease;">🧬 Drug Mechanisms</span>
            <span class="welcome-tag" style="background: #ffffff; padding: 0.5rem 1rem; border-radius: 20px; border: 1px solid #e0e7ff; color: inherit; transition: all 0.3s ease;">⚗️ Interactions</span>
            <span class="welcome-tag" style="background: #ffffff; padding: 0.5rem 1rem; border-radius: 20px; border: 1px solid #e0e7ff; color: inherit; transition: all 0.3s ease;">📊 Pharmacokinetics</span>
            <span class="welcome-tag" style="background: #ffffff; padding: 0.5rem 1rem; border-radius: 20px; border: 1px solid #e0e7ff; color: inherit; transition: all 0.3s ease;">🏥 Clinical Applications</span>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    # Quick start examples
    st.markdown("### 💡 Try these examples:")
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.button("🔬 How do ACE inhibitors work?", use_container_width=True):
            st.session_state.example_prompt = "How do ACE inhibitors work?"
            st.rerun()
        
        if st.button("💊 Warfarin drug interactions", use_container_width=True):
            st.session_state.example_prompt = "What are the major drug interactions with warfarin?"
            st.rerun()
    
    with col2:
        if st.button("⚡ Beta-blocker selectivity", use_container_width=True):
            st.session_state.example_prompt = "Explain beta-blocker selectivity and clinical implications"
            st.rerun()
        
        if st.button("🧪 NSAID side effects", use_container_width=True):
            st.session_state.example_prompt = "What are the side effects of NSAIDs and their mechanisms?"
            st.rerun()

def render_chat_messages():
    """Render chat messages with enhanced features."""
    # Ensure we have messages to display
    if not st.session_state.get('chat_messages'):
        return
    
    for i, message in enumerate(st.session_state.chat_messages):
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
            
            # Timestamps removed for cleaner mobile experience
            
            # Add regenerate button for the last assistant message
            if (message["role"] == "assistant" and 
                i == len(st.session_state.chat_messages) - 1 and 
                i > 0 and 
                not st.session_state.get('processing_input', False)):  # Don't show during processing
                
                col1, col2, col3 = st.columns([1, 1, 3])
                with col1:
                    if st.button("🔄 Regenerate", key=f"regen_{i}"):
                        regenerate_response(i)
                
                with col2:
                    if st.button("📋 Copy", key=f"copy_{i}"):
                        st.code(message["content"])
                        st.success("✅ Response copied!")

def regenerate_response(message_index):
    """Regenerate the assistant response."""
    if st.session_state.get('processing_input', False):
        st.warning("⚠️ Please wait for the current response to complete.")
        return
    
    if message_index > 0 and st.session_state.chat_messages[message_index-1]["role"] == "user":
        user_prompt = st.session_state.chat_messages[message_index-1]["content"]
        
        # Remove the current assistant response
        st.session_state.chat_messages.pop()
        
        # Process the regeneration
        process_chat_input(user_prompt)

def render_bottom_input_area():
    """Render the bottom input area with document upload."""
    
    # Add spacing to push input to bottom
    st.markdown("<br>" * 2, unsafe_allow_html=True)
    
    # Enhanced CSS for seamless input blending
    st.markdown("""
    <style>
    .stChatInput {
        position: sticky;
        bottom: 0;
        background: transparent;
        padding: 1rem 0;
        border: none;
        z-index: 999;
        margin-top: 2rem;
    }
    .main .block-container {
        padding-bottom: 120px;
    }
    .upload-area {
        background: #f8fafc;
        border: 2px dashed #d1d5db;
        border-radius: 8px;
        padding: 1rem;
        text-align: center;
        margin-bottom: 1rem;
    }
    </style>
    """, unsafe_allow_html=True)
    
    st.markdown("---")
    
    # Document upload area
    render_document_upload()
    
    # Custom chat input that matches page styling
    if st.session_state.get('processing_input', False):
        st.info("🤔 Generating response... Please wait.")
        # Show disabled input during processing
        st.text_input(
            "Message", 
            value="Generating response...", 
            disabled=True,
            label_visibility="collapsed",
            placeholder="Generating response..."
        )
    else:
        # Create columns for input and send button
        col1, col2 = st.columns([6, 1])
        
        with col1:
            prompt = st.text_input(
                "Message",
                placeholder="Ask me anything about pharmacology...",
                label_visibility="collapsed",
                key="chat_input_field"
            )
        
        with col2:
            send_button = st.button("➤", use_container_width=True, type="primary")
        
        # Handle example prompts
        if st.session_state.get('example_prompt'):
            prompt = st.session_state.example_prompt
            st.session_state.example_prompt = None
        
        # Process input when button is clicked or Enter is pressed
        if (prompt and send_button) or (prompt and prompt != st.session_state.get('last_input', '')):
            if not st.session_state.get('processing_input', False):
                st.session_state.last_input = prompt
                process_chat_input(prompt)
                # Clear the input after processing
                st.session_state.chat_input_field = ""
                st.rerun()

def render_document_upload():
    """Render enhanced document upload area."""
    with st.expander("📎 Upload Documents (PDF, DOCX, Images)", expanded=False):
        # Determine upload limits based on user type
        is_admin = st.session_state.get('username') == 'admin'
        max_files = 10 if is_admin else 3
        user_type = "Admin" if is_admin else "User"
        
        st.info(f"📷 **Image OCR**: Only text content will be extracted from images for processing. Charts, graphs, and visual elements will not be analyzed.")
        st.info(f"📊 **{user_type} Upload Limit**: You can upload up to {max_files} documents at a time.")
        
        uploaded_files = st.file_uploader(
            f"Choose up to {max_files} files to enhance your conversation",
            type=['txt', 'pdf', 'docx', 'md', 'pptx', 'png', 'jpg', 'jpeg'],
            help=f"Upload documents, presentations, or images for context-aware responses (Max: {max_files} files)",
            accept_multiple_files=True,
            key=f"doc_upload_{st.session_state.get('current_conversation_id', 'new')}"
        )
        
        if uploaded_files:
            # Check file count limit
            if len(uploaded_files) > max_files:
                st.error(f"❌ Too many files! {user_type}s can upload up to {max_files} files at a time. You selected {len(uploaded_files)} files.")
                return
            
            # Process each uploaded file
            for uploaded_file in uploaded_files:
                # Check file size limit
                file_size_mb = uploaded_file.size / (1024 * 1024)
                if file_size_mb > MAX_FILE_SIZE_MB:
                    st.error(f"❌ File too large! Maximum size allowed is {MAX_FILE_SIZE_MB}MB. File '{uploaded_file.name}' is {file_size_mb:.1f}MB.")
                    continue
                
                process_document_upload(uploaded_file)
        
        # Show current documents
        show_current_documents()

def show_current_documents():
    """Show documents in current conversation."""
    conv_id = st.session_state.get('current_conversation_id')
    if not conv_id or 'conversation_documents' not in st.session_state:
        return
    
    documents = st.session_state.conversation_documents.get(conv_id, [])
    if documents:
        st.markdown("**📚 Documents in this conversation:**")
        for i, doc in enumerate(documents):
            col1, col2 = st.columns([3, 1])
            with col1:
                st.markdown(f"• {doc['filename']} ({doc.get('file_size', 0)} bytes)")
            with col2:
                if st.button("🗑️", key=f"del_doc_{i}", help="Remove document"):
                    # Remove document and mark for cleanup
                    documents.pop(i)
                    
                    # Clean up processing flags for this document
                    file_key = f"{doc['filename']}_{doc.get('file_size', 0)}_{conv_id}"
                    if f'processed_doc_{file_key}' in st.session_state:
                        del st.session_state[f'processed_doc_{file_key}']
                    if f'processing_doc_{file_key}' in st.session_state:
                        del st.session_state[f'processing_doc_{file_key}']
                    
                    st.success(f"✅ Removed {doc['filename']}")
                    st.rerun()

def process_document_upload(uploaded_file):
    """Process uploaded document with enhanced features."""
    logger.info(f"Starting document upload process for: {uploaded_file.name}")
    
    # Check if this file has already been processed to prevent loops
    file_key = f"{uploaded_file.name}_{uploaded_file.size}_{st.session_state.get('current_conversation_id', 'new')}"
    
    if st.session_state.get(f'processed_doc_{file_key}', False):
        logger.info(f"Document already processed: {uploaded_file.name}")
        return  # Already processed, don't process again
    
    # Mark as being processed
    st.session_state[f'processing_doc_{file_key}'] = True
    
    with st.spinner(f"Processing {uploaded_file.name}..."):
        try:
            logger.info(f"Processing document content for: {uploaded_file.name}")
            
            # Process the document directly (no import needed)
            content = process_uploaded_document_multipage(uploaded_file)
            
            if content:
                logger.info(f"Document content extracted: {len(content)} characters")
                
                # Save to conversation
                success = save_document_to_conversation_multipage(uploaded_file, content)
                
                if success:
                    # Mark as successfully processed
                    st.session_state[f'processed_doc_{file_key}'] = True
                    st.session_state[f'processing_doc_{file_key}'] = False
                    
                    st.success(f"✅ Document processed: {uploaded_file.name}")
                    st.info(f"📄 Extracted {len(content)} characters of content")
                    
                    logger.info(f"Document successfully processed and saved: {uploaded_file.name}")
                    
                    # Don't call st.rerun() here - let the natural flow handle it
                else:
                    logger.error(f"Failed to save document: {uploaded_file.name}")
                    st.error("❌ Failed to save document - check logs for details")
                    st.session_state[f'processing_doc_{file_key}'] = False
            else:
                logger.error(f"Failed to extract content from document: {uploaded_file.name}")
                st.error("❌ Failed to process document - no content extracted")
                st.session_state[f'processing_doc_{file_key}'] = False
                
        except Exception as e:
            logger.error(f"Document upload error for {uploaded_file.name}: {e}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            st.error(f"❌ Error processing document: {str(e)}")
            st.session_state[f'processing_doc_{file_key}'] = False

def process_chat_input(prompt):
    """Process chat input with enhanced features and proper response visibility."""
    logger.info(f"Processing user input: {prompt[:50]}...")
    
    # Prevent duplicate processing
    if st.session_state.get('processing_input', False):
        return
    
    st.session_state.processing_input = True
    
    try:
        # Ensure we have a conversation ID
        if not st.session_state.get('current_conversation_id'):
            try:
                from utils.conversation_manager import run_async, create_new_conversation
                conversation_id = run_async(create_new_conversation())
                if conversation_id:
                    st.session_state.current_conversation_id = conversation_id
                    logger.info(f"Created new conversation: {conversation_id}")
                else:
                    raise Exception("Database conversation creation failed")
            except Exception as conv_error:
                logger.warning(f"Database conversation creation failed: {conv_error}")
                # Fallback to UUID - create local conversation
                import uuid
                conversation_id = str(uuid.uuid4())
                st.session_state.current_conversation_id = conversation_id
                
                # Add to local conversations
                from fix_user_isolation import get_secure_conversations, secure_update_conversations
                conversations = get_secure_conversations()
                conversations[conversation_id] = {
                    "title": f"New Chat {len(conversations) + 1}",
                    "messages": [],
                    "created_at": datetime.now().isoformat(),
                    "updated_at": datetime.now().isoformat(),
                    "model": st.session_state.get('selected_model_mode', 'normal'),
                    "message_count": 0
                }
                secure_update_conversations(conversations)
                logger.info(f"Created fallback local conversation: {conversation_id}")
        
        # Add user message to session state first
        user_message = {
            "role": "user", 
            "content": prompt,
            "timestamp": datetime.now().isoformat()
        }
        st.session_state.chat_messages.append(user_message)
        
        # Mark for refresh to show user message immediately
        st.session_state.force_refresh = True
        
        # Generate assistant response with streaming
        try:
            # Check if streaming is enabled
            use_streaming = st.session_state.get('use_streaming', True)
            
            # Use the simplified response generation (streaming removed for stability)
            full_response = generate_streaming_response(prompt)
            
            if not full_response or full_response.strip() == "":
                full_response = "I apologize, but I couldn't generate a response. Please try again."
            
            # Add assistant message to chat history
            assistant_message = {
                "role": "assistant", 
                "content": full_response,
                "timestamp": datetime.now().isoformat()
            }
            st.session_state.chat_messages.append(assistant_message)
            logger.info(f"Added assistant message to chat. Total messages: {len(st.session_state.chat_messages)}")
            
            # Save to database asynchronously (non-blocking)
            try:
                from utils.conversation_manager import run_async, add_message_to_current_conversation
                # Save user message
                run_async(add_message_to_current_conversation("user", prompt))
                # Save assistant message
                run_async(add_message_to_current_conversation("assistant", full_response))
                logger.info("Messages saved to database")
            except Exception as save_error:
                logger.warning(f"Failed to save to database: {save_error}")
            
            logger.info(f"Response generated and saved: {len(full_response)} characters")
            
        except Exception as e:
            logger.error(f"Error generating response: {e}")
            error_message = f"Sorry, I encountered an error: {str(e)}"
            
            # Add error message to chat history
            st.session_state.chat_messages.append({
                "role": "assistant", 
                "content": error_message,
                "timestamp": datetime.now().isoformat()
            })
    
    finally:
        # Always reset processing flag
        st.session_state.processing_input = False
        
        # Force rerun to display the response
        st.rerun()

def generate_streaming_response(prompt):
    """Generate streaming AI response for real-time display."""
    try:
        # Import AI functions
        from openai_client import chat_completion_stream, get_available_model_modes
        from prompts import pharmacology_system_prompt, pharmacology_fast_prompt
        
        # Get available models
        available_modes = get_available_model_modes()
        if not available_modes:
            return "❌ No AI models available. Please check your API keys."
        
        # Use selected model mode
        selected_mode = st.session_state.get('selected_model_mode', 'fast')  # Default to fast for speed
        if selected_mode not in available_modes:
            selected_mode = list(available_modes.keys())[0]
        
        model = available_modes[selected_mode]["model"]
        
        # Get document context (cached for speed)
        conversation_id = st.session_state.get('current_conversation_id', 'default')
        document_context = get_conversation_context_cached(prompt, conversation_id)
        
        # Choose prompt based on performance settings
        use_fast_prompt = selected_mode == 'fast'
        
        # Simplified prompt construction for speed
        if document_context and is_useful_document_context(document_context):
            from prompts import get_rag_enhanced_prompt
            enhanced_system_prompt = get_rag_enhanced_prompt(prompt, document_context)
            api_messages = [{"role": "system", "content": enhanced_system_prompt}]
        else:
            base_prompt = pharmacology_fast_prompt if use_fast_prompt else pharmacology_system_prompt
            api_messages = [
                {"role": "system", "content": base_prompt},
                {"role": "user", "content": prompt}
            ]
        
        # Use fast completion instead of streaming to avoid UI conflicts
        from openai_client import chat_completion_fast
        full_response = chat_completion_fast(model, api_messages)
        
        logger.info(f"Generated response: {len(full_response)} characters")
        return full_response
        
    except Exception as e:
        logger.error(f"Error in streaming response: {e}")
        return f"Sorry, I encountered an error: {str(e)}"

def generate_enhanced_response(prompt):
    """Generate enhanced AI response with document context (fast mode)."""
    try:
        # Import AI functions
        from openai_client import chat_completion_fast, get_available_model_modes
        from prompts import pharmacology_system_prompt, pharmacology_fast_prompt
        
        # Get available models
        available_modes = get_available_model_modes()
        if not available_modes:
            return "❌ No AI models available. Please check your API keys."
        
        # Use selected model mode
        selected_mode = st.session_state.get('selected_model_mode', 'fast')
        if selected_mode not in available_modes:
            selected_mode = list(available_modes.keys())[0]
        
        model = available_modes[selected_mode]["model"]
        
        # Get document context (cached for performance)
        conversation_id = st.session_state.get('current_conversation_id', 'default')
        document_context = get_conversation_context_cached(prompt, conversation_id)
        
        # Choose prompt based on performance settings
        use_fast_prompt = selected_mode == 'fast'
        
        # Use RAG-enhanced prompt if we have document context
        if document_context and is_useful_document_context(document_context):
            from prompts import get_rag_enhanced_prompt
            enhanced_system_prompt = get_rag_enhanced_prompt(prompt, document_context)
            
            # For RAG-enhanced prompt, we use it as the system message
            api_messages = [
                {"role": "system", "content": enhanced_system_prompt}
            ]
        else:
            # Use appropriate prompt based on speed settings
            base_prompt = pharmacology_fast_prompt if use_fast_prompt else pharmacology_system_prompt
            api_messages = [
                {"role": "system", "content": base_prompt},
                {"role": "user", "content": prompt}
            ]
        
        logger.info(f"Sending to model: {model}")
        logger.info(f"System prompt length: {len(enhanced_system_prompt)}")
        logger.info(f"User prompt: {prompt[:100]}...")
        
        # Use fast completion for maximum speed
        from openai_client import chat_completion_fast
        full_response = chat_completion_fast(model, api_messages)
        
        logger.info(f"Generated response: {len(full_response)} characters")
        return full_response
        
    except Exception as e:
        logger.error(f"Error in generate_enhanced_response: {e}")
        return f"Sorry, I encountered an error: {str(e)}"



def is_useful_document_context(context):
    """Check if document context is useful and not just error messages."""
    if not context or len(context.strip()) < 50:
        return False
    
    # Check for common error indicators
    error_indicators = [
        "OCR not available",
        "OCR failed",
        "Text extraction failed",
        "No extractable text found",
        "Processing failed",
        "install pytesseract",
        "Unsupported format",
        "text extraction is not available",
        "Please describe what you see"
    ]
    
    # If context is mostly error messages, it's not useful
    error_count = sum(1 for indicator in error_indicators if indicator.lower() in context.lower())
    
    # If more than 1 error indicator or context is very short, consider it not useful
    if error_count > 1 or len(context.strip()) < 100:
        return False
    
    return True

def clean_document_content(content):
    """Clean document content to remove unhelpful error messages."""
    if not content:
        return content
    
    # Remove common error message patterns
    error_patterns = [
        r"\[.*?OCR not available.*?\]",
        r"\[.*?OCR failed.*?\]",
        r"\[.*?Text extraction failed.*?\]",
        r"\[.*?install pytesseract.*?\]",
        r"\[.*?Processing failed.*?\]"
    ]
    
    import re
    cleaned_content = content
    for pattern in error_patterns:
        cleaned_content = re.sub(pattern, "", cleaned_content, flags=re.IGNORECASE | re.DOTALL)
    
    # Clean up extra whitespace
    cleaned_content = re.sub(r'\n\s*\n', '\n\n', cleaned_content.strip())
    
    return cleaned_content

@st.cache_data(ttl=300)  # Cache for 5 minutes
def get_conversation_context_cached(user_query, conversation_id):
    """Get document context for current conversation (cached)."""
    return get_conversation_context_uncached(user_query)

def get_conversation_context_uncached(user_query=""):
    """Get document context for current conversation - simplified version."""
    try:
        conv_id = st.session_state.get('current_conversation_id')
        if not conv_id:
            return ""
        
        # Use simple session-based documents only (avoid RAG complexity for now)
        if 'conversation_documents' in st.session_state:
            documents = st.session_state.conversation_documents.get(conv_id, [])
            if documents:
                context_parts = []
                
                for doc in documents:
                    content = doc.get('content', '')
                    filename = doc.get('filename', 'Unknown')
                    
                    # Only include documents with useful content
                    if content and len(content.strip()) > 50 and is_useful_document_context(content):
                        # Limit content length
                        clean_content = clean_document_content(content[:1500])
                        context_parts.append(f"=== {filename} ===\n{clean_content}")
                
                if context_parts:
                    return "\n\n".join(context_parts)
        
        return ""
        
    except Exception as e:
        logger.error(f"Error getting conversation context: {e}")
        return ""
    
    return loop.run_until_complete(coro)

def save_conversation_to_database():
    """Save current conversation to database."""
    try:
        if not st.session_state.get('authenticated') or not st.session_state.get('user_id'):
            logger.warning("Cannot save conversation: user not authenticated")
            return False
        
        if not st.session_state.chat_messages:
            logger.info("No messages to save")
            return True
        
        # Import services
        from services.conversation_service import conversation_service
        from services.user_service import user_service
        
        # Get user UUID
        user_data = run_async_operation(user_service.get_user_by_id(st.session_state.user_id))
        if not user_data:
            logger.error("User data not found")
            return False
        
        # Create or update conversation
        if not st.session_state.current_conversation_id:
            # Generate title from first user message
            title = generate_conversation_title()
            conversation_id = run_async_operation(conversation_service.create_conversation(
                user_data['id'], 
                title, 
                st.session_state.get('selected_model_mode', 'normal')
            ))
            st.session_state.current_conversation_id = conversation_id
            logger.info(f"Created new conversation: {conversation_id}")
        
        # Update conversation with current messages
        success = run_async_operation(conversation_service.update_conversation(
            user_data['id'],
            st.session_state.current_conversation_id,
            {'messages': st.session_state.chat_messages}
        ))
        
        if success:
            # Update local conversations cache securely
            from fix_user_isolation import get_secure_conversations, secure_update_conversation
            conv_id = st.session_state.current_conversation_id
            secure_update_conversation(conv_id, {
                'title': generate_conversation_title(),
                'messages': st.session_state.chat_messages.copy(),
                'updated_at': datetime.now().isoformat(),
                'model': st.session_state.get('selected_model_mode', 'normal'),
                'message_count': len(st.session_state.chat_messages)
            })
            
            logger.info(f"Conversation saved successfully: {len(st.session_state.chat_messages)} messages")
            return True
        else:
            logger.error("Failed to save conversation to database")
            return False
        
    except Exception as e:
        logger.error(f"Error saving conversation: {e}")
        return False

def generate_conversation_title():
    """Generate a title for the conversation based on the first user message."""
    if st.session_state.chat_messages:
        first_user_msg = next((msg for msg in st.session_state.chat_messages if msg["role"] == "user"), None)
        if first_user_msg:
            content = first_user_msg["content"]
            # Create title from first 50 characters
            title = content[:50] + "..." if len(content) > 50 else content
            return title
    
    # Fallback title
    return f"Chat {datetime.now().strftime('%m/%d %H:%M')}"

if __name__ == "__main__":
    main()